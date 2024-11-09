import os
import fitz  # PyMuPDF for PDF processing
from pptx import Presentation
import subprocess
from llama_index.core import Document
from utils import (
    describe_image, is_graph, process_graph, extract_text_around_item, 
    process_text_blocks, save_uploaded_file
)

def get_pdf_documents(pdf_file):
    """Process a PDF file and extract text, tables, and images."""
    try:
        f = fitz.open(stream=pdf_file.read(), filetype="pdf")
    except Exception as e:
        print(f"Error opening PDF file: {e}")
        return []
    
    all_pdf_documents = []
    ongoing_tables = {}

    for i in range(len(f)):
        page = f[i]
        text_blocks = [block for block in page.get_text("blocks", sort=True)
                       if block[-1] == 0 and 0.1 < block[1] / page.rect.height < 0.9]
        
        grouped_text_blocks = process_text_blocks(text_blocks)
        table_docs, table_bboxes, ongoing_tables = parse_all_tables(f, page, i, text_blocks, ongoing_tables)
        all_pdf_documents.extend(table_docs)

        image_docs = parse_all_images(f, page, i, text_blocks)
        all_pdf_documents.extend(image_docs)

        for text_block_ctr, (heading_block, content) in enumerate(grouped_text_blocks, 1):
            if not any(fitz.Rect(heading_block[:4]).intersects(bbox) for bbox in table_bboxes):
                text_doc = create_text_document(f, i, text_block_ctr, heading_block, content)
                all_pdf_documents.append(text_doc)

    f.close()
    return all_pdf_documents

def create_text_document(pdf_file, page_num, text_block_ctr, heading_block, content):
    """Helper to create a text Document with metadata."""
    bbox = {"x1": heading_block[0], "y1": heading_block[1], "x2": heading_block[2], "y2": heading_block[3]}
    return Document(
        text=f"{heading_block[4]}\n{content}",
        metadata={**bbox, "type": "text", "page_num": page_num, "source": f"{pdf_file.name[:-4]}-page{page_num}-block{text_block_ctr}"},
        id_=f"{pdf_file.name[:-4]}-page{page_num}-block{text_block_ctr}"
    )

def parse_all_tables(doc, page, page_num, text_blocks, ongoing_tables):
    """Extract tables from a PDF page."""
    table_docs, table_bboxes = [], []
    try:
        for tab in page.find_tables(horizontal_strategy="lines_strict", vertical_strategy="lines_strict"):
            pandas_df = tab.to_pandas()
            table_path = save_table_to_file(pandas_df, page_num, len(table_docs) + 1)
            before_text, after_text = extract_text_around_item(text_blocks, fitz.Rect(tab.bbox), page.rect.height)
            caption = before_text.replace("\n", " ") + process_graph(page.get_pixmap(clip=tab.bbox).tobytes()) + after_text.replace("\n", " ")
            doc_metadata = {
                "source": f"{doc.name[:-4]}-page{page_num}-table{len(table_docs) + 1}",
                "dataframe": table_path,
                "image": save_image(page.get_pixmap(clip=tab.bbox), page_num, len(table_docs) + 1),
                "caption": caption,
                "type": "table",
                "page_num": page_num
            }
            table_docs.append(Document(text=f"This is a table with caption: {caption}", metadata=doc_metadata))
            table_bboxes.append(fitz.Rect(tab.bbox))
    except Exception as e:
        print(f"Error extracting tables: {e}")
    return table_docs, table_bboxes, ongoing_tables

def parse_all_images(doc, page, page_num, text_blocks):
    """Extract images from a PDF page."""
    image_docs = []
    for image_info in page.get_image_info(xrefs=True):
        xref = image_info['xref']
        img_bbox = fitz.Rect(image_info['bbox'])
        if xref and valid_image_size(img_bbox, page.rect):
            img_data = page.parent.extract_image(xref)["image"]
            before_text, after_text = extract_text_around_item(text_blocks, img_bbox, page.rect.height)
            caption = before_text.replace("\n", " ") + (process_graph(img_data) if is_graph(img_data) else "") + after_text.replace("\n", " ")
            image_docs.append(Document(
                text="This is an image with caption: " + caption,
                metadata={"source": f"{doc.name[:-4]}-page{page_num}-image{xref}", "image": save_image_data(img_data, page_num, xref), "caption": caption, "type": "image", "page_num": page_num}
            ))
    return image_docs

def valid_image_size(bbox, page_rect):
    """Check if image size is reasonable."""
    return bbox.width > page_rect.width / 20 and bbox.height > page_rect.height / 20

def save_table_to_file(dataframe, page_num, table_num):
    """Save table to an Excel file."""
    table_path = os.path.join("vectorstore/table_references", f"table{table_num}-page{page_num}.xlsx")
    os.makedirs(os.path.dirname(table_path), exist_ok=True)
    dataframe.to_excel(table_path)
    return table_path

def save_image(pixmap, page_num, img_num):
    """Save image from pixmap."""
    image_path = os.path.join("vectorstore/image_references", f"image{img_num}-page{page_num}.jpg")
    os.makedirs(os.path.dirname(image_path), exist_ok=True)
    pixmap.save(image_path)
    return image_path

def save_image_data(img_data, page_num, img_num):
    """Save image data to a file."""
    image_path = os.path.join("vectorstore/image_references", f"image{img_num}-page{page_num}.png")
    os.makedirs(os.path.dirname(image_path), exist_ok=True)
    with open(image_path, "wb") as img_file:
        img_file.write(img_data)
    return image_path

def process_ppt_file(ppt_path):
    """Process a PowerPoint file."""
    pdf_path = convert_ppt_to_pdf(ppt_path)
    images_data = convert_pdf_to_images(pdf_path)
    slide_texts = extract_text_and_notes_from_ppt(ppt_path)
    processed_data = []
    for (image_path, page_num), (slide_text, notes) in zip(images_data, slide_texts):
        processed_data.append(Document(
            text="This is a slide with text: " + slide_text + (process_graph(slide_text.encode()) if is_graph(slide_text.encode()) else ""),
            metadata={"source": ppt_path, "image": image_path, "caption": slide_text + notes, "type": "image", "page_num": page_num}
        ))
    return processed_data

def convert_ppt_to_pdf(ppt_path):
    """Convert PPT to PDF with LibreOffice."""
    pdf_path = os.path.join("vectorstore/ppt_references", f"{os.path.splitext(os.path.basename(ppt_path))[0]}.pdf")
    os.makedirs(os.path.dirname(pdf_path), exist_ok=True)
    subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', os.path.dirname(pdf_path), ppt_path], check=True)
    return pdf_path

def convert_pdf_to_images(pdf_path):
    """Convert PDF to images."""
    doc = fitz.open(pdf_path)
    image_paths = [(save_image(fitz.Pixmap(doc.load_page(i).get_pixmap()), i, i), i) for i in range(len(doc))]
    doc.close()
    return image_paths

def extract_text_and_notes_from_ppt(ppt_path):
    """Extract text and notes from a PowerPoint file."""
    prs = Presentation(ppt_path)
    return [( ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")]), slide.notes_slide.notes_text_frame.text if slide.notes_slide else '') for slide in prs.slides]

def load_multimodal_data(files):
    """Load and process multiple file types."""
    documents = []
    for file in files:
        try:
            file_extension = os.path.splitext(file.name.lower())[1]
            if file_extension in ('.png', '.jpg', '.jpeg'):
                documents.append(Document(text=describe_image(file.read()), metadata={"source": file.name, "type": "image"}))
            elif file_extension == '.pdf':
                documents.extend(get_pdf_documents(file))
            elif file_extension in ('.ppt', '.pptx'):
                documents.extend(process_ppt_file(save_uploaded_file(file)))
            else:
                documents.append(Document(text=file.read().decode("utf-8"), metadata={"source": file.name, "type": "text"}))
        except Exception as e:
            print(f"Error processing {file.name}: {e}")
    return documents

def load_data_from_directory(directory):
    """Load and process files from a directory."""
    return load_multimodal_data([open(os.path.join(directory, filename), "rb") for filename in os.listdir(directory) if os.path.isfile(os.path.join(directory, filename))])